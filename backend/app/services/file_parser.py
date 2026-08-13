import os
import tempfile
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".txt"}


class TextExtractionError(Exception):
    pass


def extract_text(filepath: str, original_name: str) -> str:
    ext = os.path.splitext(original_name)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                pages = [
                    f"[Page {i + 1}]\n{(p.extract_text() or '')}"
                    for i, p in enumerate(reader.pages)
                ]
                text = "\n\n".join(pages)
        elif ext == ".docx":
            doc = DocxDocument(filepath)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(c for c in cells if c))
            text = "\n\n".join(parts)
        elif ext == ".pptx":
            prs = Presentation(filepath)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                lines = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
                if lines:
                    slides.append(f"[Slide {i}]\n" + "\n".join(lines))
            text = "\n\n".join(slides)
        elif ext in (".xlsx", ".xls"):
            xf = pd.ExcelFile(filepath)
            parts = [
                f"[Sheet: {s}]\n{pd.read_excel(xf, sheet_name=s).to_string(index=False)}"
                for s in xf.sheet_names
            ]
            text = "\n\n".join(parts)
        elif ext == ".csv":
            text = f"[CSV Data]\n{pd.read_csv(filepath, on_bad_lines='skip').to_string(index=False)}"
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
    except Exception as e:
        raise TextExtractionError(
            f"Could not extract text from {original_name}: {e}"
        ) from e
    return text.strip()


def extract_text_from_bytes(data: bytes, original_name: str) -> str:
    ext = os.path.splitext(original_name)[1].lower()
    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(data)
        tmp_path = tmp.name
    try:
        return extract_text(tmp_path, original_name)
    finally:
        os.unlink(tmp_path)
